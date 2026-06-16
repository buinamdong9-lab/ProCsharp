from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(56, 189, 248)
WHITE = RGBColor(240, 248, 255)
MUTED = RGBColor(191, 219, 254)


def add_title_slide(title, subtitle, note=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    bar = slide.shapes.add_shape(1, Inches(0.35), Inches(0.35), Inches(0.18), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.color.rgb = BLUE

    tx = slide.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(11.5), Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    tx2 = slide.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(11.0), Inches(1.0))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED

    if note:
        tx3 = slide.shapes.add_textbox(Inches(0.95), Inches(6.1), Inches(11.2), Inches(0.5))
        tf3 = tx3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(148, 163, 184)


def add_content_slide(title, bullets, footer=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    bar = slide.shapes.add_shape(1, Inches(0.35), Inches(0.35), Inches(0.18), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.2), Inches(4.8))
    tfb = body.text_frame
    tfb.word_wrap = True
    for i, item in enumerate(bullets):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    if footer:
        foot = slide.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(11.2), Inches(0.35))
        tfc = foot.text_frame
        pc = tfc.paragraphs[0]
        pc.text = footer
        pc.font.size = Pt(10)
        pc.font.color.rgb = RGBColor(191, 219, 254)


add_title_slide('Hệ thống quản lý mượn - trả thiết bị', 'Bản trình chiếu tổng quan về giải pháp, chức năng và quy trình vận hành.', 'Đồ án C# - Bài tập lớn')
add_content_slide('Vấn đề cần giải quyết', [
    'Quản lý nhiều thiết bị, phòng học và người dùng phân tán khó kiểm soát.',
    'Phiếu mượn/trả thiếu minh bạch, dễ phát sinh sai sót hoặc bỏ sót.',
    'Cần quy trình phê duyệt, theo dõi trạng thái và báo cáo nhanh.'
], 'Slide 2 / 5')
add_content_slide('Chức năng chính', [
    'Quản lý thiết bị, phòng học, người dùng và phiếu mượn/trả.',
    'Phê duyệt yêu cầu, cập nhật trạng thái và hỗ trợ tìm kiếm, lọc dữ liệu.',
    'Quản lý thùng rác, cấu hình hệ thống và giao diện WinForms thân thiện.'
], 'Slide 3 / 5')
add_content_slide('Luồng xử lý', [
    'Người dùng đăng nhập và tạo yêu cầu mượn thiết bị.',
    'Hệ thống kiểm tra thiết bị, phòng học và trạng thái trước khi duyệt.',
    'Quản lý theo dõi việc phê duyệt, cập nhật trả thiết bị và báo cáo.'
], 'Slide 4 / 5')
add_content_slide('Kết luận', [
    'Giúp tăng tính minh bạch, giảm sai sót và quản lý thiết bị hiệu quả hơn.',
    'Là nền tảng phù hợp cho việc mở rộng báo cáo, phân quyền và tích hợp dữ liệu sau này.'
], 'Slide 5 / 5')

prs.save('BaoCao/rendered/presentation.pptx')
print('Saved BaoCao/rendered/presentation.pptx')
